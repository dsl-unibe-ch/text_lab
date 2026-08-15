"""
Format-preserving document translation.

For each supported document type this module returns bytes of the **same**
file type with the natural-language text translated in place, leaving all
structural markup (headings, lists, tables, images, links, code, math)
intact wherever possible.

Public API
----------

    translate_markdown(md_text, translate_fn, progress_cb=None) -> str
    translate_docx(docx_bytes, translate_fn, progress_cb=None) -> bytes
    translate_pdf(pdf_bytes, translate_fn, progress_cb=None,
                  target_iso=None) -> bytes

Every ``translate_fn`` is a plain callable ``str -> str`` — usually a
partially-applied :func:`core.translation.engine.translate`. The functions in
this module handle:

* **Markdown** — line-by-line prose translation with structure detection
  (headings, lists, blockquotes) plus :mod:`core.translation.shield` masking to
  protect links, math, code, and inline HTML.
* **DOCX** — raw XML walk (stdlib only) of ``word/document.xml`` (plus any
  headers/footers). Paragraph runs are translated in place while keeping
  paragraph styles, run styles, tables, images, and lists intact.
* **PDF** — pymupdf block extraction, semantic reflow (hyphen joining, line
  merging, cross-block paragraph detection), redaction of the original text
  bounding boxes preserving images, and reinsertion of the translated text
  at the same location with auto-shrinking font size.

The module is Streamlit-free. All progress reporting is via ``progress_cb``
callbacks with signature ``(done: int, total: int, stage: str)``.
"""

from __future__ import annotations

import functools
import io
import os
import re
import tempfile
import zipfile
from pathlib import Path
from typing import Callable, Dict, List, Mapping, Optional, Tuple
from xml.etree import ElementTree as ET

from .shield import shielded_translate, shielded_translate_many

TranslateFn = Callable[[str], str]
ProgressCb = Optional[Callable[[int, int, str], None]]
Glossary = Optional[Mapping[str, str]]

# Scanned-PDF threshold: total non-whitespace characters in the text layer
# below which the PDF is treated as image-only and routed through OCR.
SCANNED_PDF_TEXT_THRESHOLD = 200


def _report(cb: ProgressCb, done: int, total: int, stage: str) -> None:
    if cb is not None:
        try:
            cb(done, total, stage)
        except Exception:
            pass


# ===========================================================================
# Markdown
# ===========================================================================

# Structural line matchers. We only translate the *prose* portion of a line
# and re-attach the structural prefix afterwards.
_MD_HEADING_RE = re.compile(r"^(\s{0,3}#{1,6}\s+)(.*?)(\s+#+\s*)?$")
_MD_LIST_RE = re.compile(r"^(\s*(?:[-*+]|\d+[\.\)])\s+(?:\[[ xX]\]\s+)?)(.+)$")
_MD_BLOCKQUOTE_RE = re.compile(r"^(\s*>+\s*)(.*)$")
_MD_FENCE_RE = re.compile(r"^\s*(?:```|~~~)")
_MD_HR_RE = re.compile(r"^\s*(?:-{3,}|_{3,}|\*{3,})\s*$")


def translate_markdown(
    md_text: str,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
) -> str:
    """
    Translate a Markdown document while preserving structure.

    Fenced code blocks are passed through untouched. Every other line is
    routed through :func:`shielded_translate` so links, math, inline code,
    HTML, and placeholders survive the round-trip. The optional
    ``glossary`` maps source-language terms to forced target-language
    replacements.

    All translatable line bodies are collected first and translated in a
    single batched call (:func:`shielded_translate_many`), which on GPU is
    dramatically faster than one model call per line.
    """
    lines = md_text.splitlines(keepends=False)
    out: List[Optional[str]] = []
    in_fence = False
    total = len(lines)

    # Gather translatable bodies; each slot records where/how to reinsert.
    bodies: List[str] = []
    slots: List[Tuple[int, str, str]] = []  # (out_index, prefix, suffix)

    def _defer(prefix: str, body: str, suffix: str) -> None:
        out.append(None)
        slots.append((len(out) - 1, prefix, suffix))
        bodies.append(body)

    for i, line in enumerate(lines, start=1):
        _report(progress_cb, i, total, "parsing markdown")

        # Code fences: toggle and passthrough (fence + contents).
        if _MD_FENCE_RE.match(line):
            in_fence = not in_fence
            out.append(line)
            continue
        if in_fence:
            out.append(line)
            continue

        # Blank line / HR / structural-only lines: passthrough.
        if not line.strip() or _MD_HR_RE.match(line):
            out.append(line)
            continue

        # Heading:  ## Title
        m = _MD_HEADING_RE.match(line)
        if m:
            prefix, body, suffix = m.group(1), m.group(2), (m.group(3) or "")
            _defer(prefix, body, suffix)
            continue

        # Blockquote:  > text
        m = _MD_BLOCKQUOTE_RE.match(line)
        if m:
            prefix, body = m.group(1), m.group(2)
            if body.strip():
                _defer(prefix, body, "")
            else:
                out.append(line)
            continue

        # List item:  - text  |  1. text  |  * [x] text
        m = _MD_LIST_RE.match(line)
        if m:
            prefix, body = m.group(1), m.group(2)
            _defer(prefix, body, "")
            continue

        # Regular paragraph line.
        _defer("", line, "")

    _report(progress_cb, total, total, "translating markdown")
    translated = shielded_translate_many(bodies, translate_fn, glossary=glossary)
    for (idx, prefix, suffix), tr in zip(slots, translated):
        out[idx] = f"{prefix}{tr}{suffix}"

    return "\n".join(s if s is not None else "" for s in out)


# ===========================================================================
# DOCX (raw XML — no python-docx dependency)
# ===========================================================================

_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_XML_NS = "http://www.w3.org/XML/1998/namespace"

# ElementTree namespace registration so output XML keeps the ``w:`` prefix.
ET.register_namespace("w", _W_NS)


def _q(tag: str) -> str:
    """Return the Clark-notation qualified name for a WordProcessingML tag."""
    return f"{{{_W_NS}}}{tag}"


def _iter_paragraphs(root: ET.Element):
    """Yield every ``<w:p>`` element (paragraphs) inside the given root, in
    document order. Tables are covered because they contain ``<w:p>`` too."""
    yield from root.iter(_q("p"))


def _paragraph_text(p: ET.Element) -> str:
    """Concatenate all ``<w:t>`` children of a paragraph, in order."""
    parts: List[str] = []
    for t in p.iter(_q("t")):
        if t.text:
            parts.append(t.text)
    return "".join(parts)


def _rewrite_paragraph(p: ET.Element, translated: str) -> None:
    """
    Put ``translated`` back into the paragraph.

    Strategy: place the whole translated string in the first ``<w:t>`` and
    empty the rest, preserving each run's parent ``<w:r>`` (and therefore
    its style — bold/italic/color of the first run wins for the paragraph).
    Also sets ``xml:space="preserve"`` so leading/trailing whitespace is
    kept by Word.
    """
    ts = list(p.iter(_q("t")))
    if not ts:
        return
    first, rest = ts[0], ts[1:]
    first.text = translated
    first.set(f"{{{_XML_NS}}}space", "preserve")
    for t in rest:
        t.text = ""


def _translate_docx_part(
    xml_bytes: bytes,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb,
    stage: str,
    counter: List[int],
    total_est: int,
    glossary: Glossary = None,
) -> bytes:
    """Translate a single DOCX XML part and return new bytes."""
    if not xml_bytes.strip():
        return xml_bytes
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return xml_bytes

    paragraphs = list(_iter_paragraphs(root))
    texts = [_paragraph_text(p) for p in paragraphs]
    counter[0] += len(paragraphs)
    _report(progress_cb, counter[0], total_est, stage)

    translated_list = shielded_translate_many(texts, translate_fn, glossary=glossary)
    for p, text, translated in zip(paragraphs, texts, translated_list):
        if not text.strip():
            continue
        _rewrite_paragraph(p, translated)

    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


# Parts of a .docx zip that contain user-visible prose.
_DOCX_TRANSLATABLE_PARTS = (
    "word/document.xml",
    "word/footnotes.xml",
    "word/endnotes.xml",
    "word/comments.xml",
)
_DOCX_TRANSLATABLE_PREFIXES = ("word/header", "word/footer")


def _is_translatable_docx_part(name: str) -> bool:
    if name in _DOCX_TRANSLATABLE_PARTS:
        return True
    return name.endswith(".xml") and any(
        name.startswith(pfx) for pfx in _DOCX_TRANSLATABLE_PREFIXES
    )


def translate_docx(
    docx_bytes: bytes,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
) -> bytes:
    """
    Translate a .docx file, returning a new .docx file.

    Preserves: paragraph styles (headings, list styles), tables (each cell
    is a paragraph in the XML), inline images, hyperlinks, footnotes/
    endnotes, comments, headers and footers.

    Best-effort limitations:
    * If a paragraph is split into multiple styled runs (e.g. a bold word in
      the middle of a sentence), the whole paragraph is translated as one
      unit and re-inserted into the first run. Intra-paragraph run styling
      is therefore lost; paragraph-level style is preserved.
    """
    src = zipfile.ZipFile(io.BytesIO(docx_bytes), "r")
    out_buf = io.BytesIO()
    dst = zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED)

    # First pass: estimate paragraph count for progress reporting.
    total_paragraphs = 0
    parts_to_translate = []
    for info in src.infolist():
        if _is_translatable_docx_part(info.filename):
            data = src.read(info.filename)
            try:
                root = ET.fromstring(data)
                total_paragraphs += sum(1 for _ in _iter_paragraphs(root))
            except ET.ParseError:
                pass
            parts_to_translate.append(info.filename)

    total_paragraphs = max(1, total_paragraphs)
    counter = [0]

    for info in src.infolist():
        data = src.read(info.filename)
        if info.filename in parts_to_translate:
            data = _translate_docx_part(
                data,
                translate_fn,
                progress_cb,
                "translating docx",
                counter,
                total_paragraphs,
                glossary=glossary,
            )
        dst.writestr(info, data)

    src.close()
    dst.close()
    return out_buf.getvalue()


# ===========================================================================
# PDF (pymupdf)
# ===========================================================================

_ENDS_SENTENCE_RE = re.compile(r"[\.\!\?\u3002\uFF01\uFF1F\u203C\u2049\uFF0E]\s*[\"'\)\]]?\s*$")


# --- Math detection --------------------------------------------------------
#
# Research PDFs commonly render equations either in a math-specific font
# (Computer Modern math families, STIX Math, Cambria Math, MTMain,
# Latin Modern Math, ...) or as high-density blocks of math operators.
# Sending such content through a NMT model destroys the equation and
# often produces garbled fallback glyphs after re-insertion because the
# target font lacks the original math symbols. We therefore detect
# "math-y" blocks and leave them untouched: no translation, no redaction.
#
# IMPORTANT: the font pattern is deliberately narrow. Earlier versions
# matched "CMR" (Computer Modern Roman) and "Latin Modern" — the body
# fonts of most LaTeX papers — which caused every text block to be
# flagged as math, producing an unchanged output document.
_MATH_FONT_PATTERN = re.compile(
    r"\bCMSY\d*"                       # Computer Modern Symbols
    r"|\bCMMI\d*"                       # Computer Modern Math Italic
    r"|\bCMEX\d*"                       # Computer Modern Math Extension
    r"|\bMSAM\d*|\bMSBM\d*"             # AMS math fonts
    r"|\bMTMain\b|\bMTSym\b|\bMTEx\b|\bMTExtra\b"       # MathTime
    r"|LatinModern-?Math\b"             # ONLY the Math variant (not Roman/Mono/Italic)
    r"|Cambria\s?Math\b"                # OpenType math
    r"|(?:STIX|STIXTwo)\s?Math\b"       # STIX Math families only
    r"|MathJax"
    r"|Lucida(?:New)?Math\b"
    r"|\bSymbol(?:MT)?\b"               # Adobe Symbol / SymbolMT
    r"|Euler(?:Fraktur|Script|Extra)?\b",
    re.IGNORECASE,
)


def _is_math_char(ch: str) -> bool:
    """Return True for characters that are almost exclusively used in equations."""
    code = ord(ch)
    # NOTE: intentionally excludes the Greek block (0x0370-0x03FF) and the
    # Arrows block (0x2190-0x21FF). Both appear in body text often enough
    # (proper nouns, bullets, decorative arrows) to make them unreliable.
    return (
        0x2200 <= code <= 0x22FF        # Math Operators (strong)
        or 0x2A00 <= code <= 0x2AFF     # Supplemental Math Operators
        or 0x27C0 <= code <= 0x27EF     # Miscellaneous Math Symbols-A
        or 0x2980 <= code <= 0x29FF     # Miscellaneous Math Symbols-B
        or 0x1D400 <= code <= 0x1D7FF   # Math Alphanumeric Symbols
        or 0x2100 <= code <= 0x214F     # Letterlike Symbols (ℕ ℤ ℝ ∑ ...)
    )


def _is_math_span(span: dict) -> bool:
    """Heuristic: does a pymupdf span look like part of an equation?"""
    font = str(span.get("font", ""))
    if font and _MATH_FONT_PATTERN.search(font):
        return True
    text = (span.get("text", "") or "").strip()
    if len(text) < 3:
        # Ignore tiny spans (bullets, subscripts, single punctuation).
        return False
    math_count = sum(1 for c in text if _is_math_char(c))
    if math_count == 0:
        return False
    # Require a strong ratio *and* at least two math characters so isolated
    # symbols in a sentence never flip a whole span.
    return math_count >= 2 and math_count / len(text) >= 0.30


def _is_math_block(block: dict) -> bool:
    """
    Return True if a block should be treated as an equation / formula and
    left untouched. Requires a strict majority of spans to look like math.
    """
    spans = [s for line in block.get("lines", []) for s in line.get("spans", [])]
    if not spans:
        return False
    math_spans = sum(1 for s in spans if _is_math_span(s))
    return math_spans * 2 > len(spans)  # STRICT > 50%


# --- Image-overlap detection ----------------------------------------------


def _rect_area(rect: Tuple[float, float, float, float]) -> float:
    return max(0.0, rect[2] - rect[0]) * max(0.0, rect[3] - rect[1])


def _rect_intersect_area(
    a: Tuple[float, float, float, float],
    b: Tuple[float, float, float, float],
) -> float:
    ix0 = max(a[0], b[0])
    iy0 = max(a[1], b[1])
    ix1 = min(a[2], b[2])
    iy1 = min(a[3], b[3])
    return max(0.0, ix1 - ix0) * max(0.0, iy1 - iy0)


def _extract_image_bboxes(page_dict: dict) -> List[Tuple[float, float, float, float]]:
    """Return the list of image-block bboxes on a page (block ``type == 1``)."""
    out: List[Tuple[float, float, float, float]] = []
    for block in page_dict.get("blocks", []):
        if block.get("type") == 1:
            bbox = block.get("bbox")
            if bbox and len(bbox) == 4:
                out.append(tuple(bbox))
    return out


def _text_block_overlaps_image(
    text_bbox: Tuple[float, float, float, float],
    image_bboxes: List[Tuple[float, float, float, float]],
    min_overlap: float = 0.40,
) -> bool:
    """
    Return True if the text block's bbox is significantly enclosed by an
    image bbox (default 40% of the text-block area).

    Redacting such blocks tends to leave white bands over figures and
    corrupt visual context, so we skip translation for them.
    """
    ta = _rect_area(text_bbox)
    if ta <= 0:
        return False
    for ib in image_bboxes:
        if _rect_intersect_area(text_bbox, ib) / ta >= min_overlap:
            return True
    return False


# --- Unicode font selection for text re-insertion -------------------------
#
# Helvetica (pymupdf's "helv") is a Base-14 font limited to WinAnsi. It
# cannot render most Greek/math symbols, CJK, Arabic, Devanagari, etc.
# When available we register a system TrueType font and use it instead.
_UNICODE_FONT_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
)


@functools.lru_cache(maxsize=1)
def _unicode_font_path() -> Optional[str]:
    """Return the first available Unicode TrueType font, or ``None``."""
    for path in _UNICODE_FONT_CANDIDATES:
        if os.path.isfile(path):
            return path
    return None


def _reflow_lines(lines: List[str]) -> str:
    """
    Join lines belonging to the same paragraph.

    * Trailing hyphen at end of a line + next line starts lowercase =>
      remove hyphen and join with no space.
    * Otherwise join with a single space, collapsing whitespace.
    """
    out_parts: List[str] = []
    for i, raw in enumerate(lines):
        line = raw.rstrip()
        if not line:
            continue
        if out_parts and out_parts[-1].endswith("-") and line[:1].islower():
            out_parts[-1] = out_parts[-1][:-1] + line
        else:
            if out_parts and not out_parts[-1].endswith(" "):
                out_parts.append(" ")
            out_parts.append(line)
    return "".join(out_parts).strip()


def _block_median_fontsize(block: dict) -> float:
    """Return a reasonable font size to use when inserting translated text."""
    sizes: List[float] = []
    for line in block.get("lines", []):
        for span in line.get("spans", []):
            s = span.get("size")
            if s:
                sizes.append(float(s))
    if not sizes:
        return 10.0
    sizes.sort()
    return sizes[len(sizes) // 2]


def _semantic_paragraphs(
    all_blocks: List[Tuple[int, tuple, str, float]],
) -> List[List[int]]:
    """
    Group block indices into semantic paragraphs.

    Two adjacent blocks belong to the same paragraph when the first block
    does not end with sentence-terminating punctuation. This handles column
    breaks and page breaks that cut a paragraph in half.
    """
    groups: List[List[int]] = []
    current: List[int] = []
    for idx, (_page, _bbox, text, _size) in enumerate(all_blocks):
        if not text.strip():
            if current:
                groups.append(current)
                current = []
            continue
        if not current:
            current = [idx]
            continue
        prev_text = all_blocks[current[-1]][2]
        if _ENDS_SENTENCE_RE.search(prev_text):
            groups.append(current)
            current = [idx]
        else:
            current.append(idx)
    if current:
        groups.append(current)
    return groups


def translate_pdf(
    pdf_bytes: bytes,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
) -> bytes:
    """
    Translate a PDF, returning a new PDF with the same page layout, images,
    and vector graphics but with the natural-language text translated.

    Strategy
    --------
    1. Extract all text blocks with bounding boxes (pymupdf ``get_text('dict')``).
    2. Reflow lines within each block (hyphen joining, whitespace).
    3. Skip blocks that look like equations (font-based + Unicode-block
       heuristics) or that are heavily enclosed by an image bbox
       (figure captions embedded in images, chart legends, ...). Those
       blocks are left completely untouched.
    4. Group the surviving blocks into *semantic paragraphs* across page
       and column boundaries so the model always sees full logical units.
    5. Translate each paragraph through :func:`shielded_translate`.
    6. Redact only the original text bboxes we translated; keep images and
       vector line-art untouched.
    7. Insert the translated text back into the same bboxes with an
       auto-shrunk font size, using a system Unicode font when available.

    Limitations
    -----------
    * Complex equations that pymupdf mis-classifies as ordinary text still
      slip through. The math detector is intentionally conservative.
    * When no Unicode TrueType font is installed we fall back to
      Helvetica (Base-14, WinAnsi only). Non-Latin scripts may then
      render as boxes.
    * Line-wrapping inside the original bbox is recomputed by pymupdf,
      so per-line breaks do not match the original layout.
    """
    import fitz  # pymupdf, already in the container

    src = fitz.open(stream=pdf_bytes, filetype="pdf")

    # ---- 1 & 2: gather every candidate text block plus its math / image
    #             flags. We DO NOT drop anything yet — we first want a
    #             document-wide view so we can disable a heuristic that
    #             mis-classifies almost everything (safety net against
    #             a false-positive rate near 100%, which used to happen
    #             when a body font was mistaken for a math font and
    #             produced an unchanged output document).
    candidates: List[dict] = []
    # Each entry: {page, bbox, reflowed, size, is_math, in_image}
    images_per_page: dict[int, List[Tuple[float, float, float, float]]] = {}

    for page_index, page in enumerate(src):
        d = page.get_text("dict")
        image_bboxes = _extract_image_bboxes(d)
        images_per_page[page_index] = image_bboxes

        for block in d.get("blocks", []):
            if block.get("type", 0) != 0:  # 0 = text; 1 = image
                continue

            line_texts = [
                "".join(sp.get("text", "") for sp in line.get("spans", []))
                for line in block.get("lines", [])
            ]
            reflowed = _reflow_lines(line_texts)
            if not reflowed.strip():
                continue

            bbox = tuple(block.get("bbox", (0, 0, 0, 0)))
            candidates.append({
                "page": page_index,
                "bbox": bbox,
                "text": reflowed,
                "size": _block_median_fontsize(block),
                "is_math": _is_math_block(block),
                "in_image": _text_block_overlaps_image(bbox, image_bboxes),
            })

    if not candidates:
        src.close()
        return pdf_bytes

    # Safety nets: if a heuristic flags an overwhelming majority of blocks,
    # something is wrong with the heuristic on this particular document —
    # disable it rather than return an unchanged file.
    math_ratio = sum(1 for c in candidates if c["is_math"]) / len(candidates)
    if math_ratio > 0.80:
        for c in candidates:
            c["is_math"] = False
    image_ratio = sum(1 for c in candidates if c["in_image"]) / len(candidates)
    if image_ratio > 0.80:
        for c in candidates:
            c["in_image"] = False

    # Apply filters.
    kept = [c for c in candidates if not c["is_math"] and not c["in_image"]]

    # Final safety net: if all candidates were filtered out, translate
    # everything anyway. Better a slightly imperfect translation than an
    # untranslated file.
    if not kept:
        kept = candidates

    all_blocks: List[Tuple[int, tuple, str, float]] = [
        (c["page"], c["bbox"], c["text"], c["size"]) for c in kept
    ]

    if not all_blocks:
        src.close()
        return pdf_bytes

    # ---- 3: semantic grouping across blocks & pages ---------------------
    groups = _semantic_paragraphs(all_blocks)
    total = len(groups)

    # ---- 4: translate all semantic paragraphs in one batched pass -------
    _report(progress_cb, 0, total, "translating pdf")
    joined_list = [" ".join(all_blocks[i][2] for i in group) for group in groups]
    translated_list = shielded_translate_many(
        joined_list, translate_fn, glossary=glossary
    )
    _report(progress_cb, total, total, "translating pdf")

    translated_per_block: dict[int, str] = {}
    for group, translated_joined in zip(groups, translated_list):
        # Distribute the translated string across the group's blocks by
        # length ratio so text stays near its original position when a
        # paragraph spans multiple blocks.
        if len(group) == 1:
            translated_per_block[group[0]] = translated_joined
            continue

        total_len = sum(len(all_blocks[i][2]) for i in group)
        pieces = _split_by_ratio(
            translated_joined,
            [len(all_blocks[i][2]) / max(1, total_len) for i in group],
        )
        for i, piece in zip(group, pieces):
            translated_per_block[i] = piece

    # ---- 5: redact originals, insert translations -----------------------
    per_page: dict[int, List[int]] = {}
    for i, (pg, _bbox, _t, _s) in enumerate(all_blocks):
        per_page.setdefault(pg, []).append(i)

    # apply_redactions constants (fall back to numeric defaults if the
    # installed pymupdf is old and doesn't expose them).
    IMG_NONE = getattr(fitz, "PDF_REDACT_IMAGE_NONE", 0)
    LINE_ART_NONE = getattr(fitz, "PDF_REDACT_LINE_ART_NONE", 0)

    font_path = _unicode_font_path()

    for page_index, page in enumerate(src):
        idxs = per_page.get(page_index, [])
        if not idxs:
            continue

        for i in idxs:
            bbox = fitz.Rect(all_blocks[i][1])
            page.add_redact_annot(bbox, fill=(1, 1, 1))

        # Explicitly keep images AND vector line art (equation frames,
        # chart axes, table rules) untouched — the crucial fix for
        # research PDFs where redaction was clobbering figures.
        try:
            page.apply_redactions(images=IMG_NONE, graphics=LINE_ART_NONE)
        except TypeError:
            # Older pymupdf without the ``graphics`` kwarg.
            page.apply_redactions(images=IMG_NONE)

        # Insert translated text in each bbox with auto-shrink.
        for i in idxs:
            bbox = fitz.Rect(all_blocks[i][1])
            text = translated_per_block.get(i, "")
            if not text.strip():
                continue
            _insert_autoshrink(page, bbox, text, all_blocks[i][3], font_path)

    out_buf = io.BytesIO()
    src.save(out_buf, garbage=3, deflate=True)
    src.close()
    return out_buf.getvalue()


# ===========================================================================
# PDF -> Markdown (uses the OCR pipeline for structure recovery and scans)
# ===========================================================================


def detect_pdf_is_scanned(
    pdf_bytes: bytes,
    threshold: int = SCANNED_PDF_TEXT_THRESHOLD,
) -> bool:
    """True when the PDF has essentially no extractable text layer.

    Cheap PyMuPDF-only check: sums stripped text length across pages and
    stops as soon as the threshold is crossed.
    """
    import fitz

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        total = 0
        for page in doc:
            total += len(page.get_text("text").strip())
            if total >= threshold:
                return False
        return True
    finally:
        doc.close()


def pdf_needs_ocr(pdf_bytes: bytes) -> bool:
    """True if translating this PDF to Markdown would invoke the OCR (VL) lane.

    Mirrors :func:`core.auto_ocr.process_document` per-page routing under the
    ``auto`` lane: a page is sent to PaddleOCR-VL only when it lacks a real
    text layer *or* contains math. A plain born-digital PDF (text layer, no
    equations) returns ``False`` and never touches the OCR worker.

    Used to decide whether a PDF can be handled on a GPU where OCR and the
    translation model can't be resident at once (e.g. a 24 GB RTX 4090).
    """
    import fitz

    from core.auto_ocr import _page_has_math, _page_has_text_layer

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for page in doc:
            if not _page_has_text_layer(page) or _page_has_math(page):
                return True
        return False
    finally:
        doc.close()


def _ocr_progress_bridge(progress_cb: ProgressCb, stage_hint: str):
    """Adapt auto_ocr's ``callback(fraction, text)`` to our ``(done,total,stage)``."""
    if progress_cb is None:
        return None

    def _cb(frac: float, text: str) -> None:
        pct = int(round(max(0.0, min(1.0, frac)) * 1000))
        # Prepend the outer stage hint so the UI shows "OCR: parsing page 3/12".
        try:
            progress_cb(pct, 1000, f"{stage_hint}: {text}")
        except Exception:
            pass

    return _cb


def pdf_to_markdown_bundle(
    pdf_bytes: bytes,
    *,
    pdf_type: str = "auto",
    source_name: str = "input.pdf",
    progress_cb: ProgressCb = None,
    free_translation_vram_first: bool = False,
) -> Tuple[str, Dict[str, bytes]]:
    """Extract ``pdf_bytes`` into (markdown_text, {asset_name: png_bytes}).

    ``pdf_type``:
      * ``"auto"``      — ``auto_ocr.process_document(native_fast_lane=True)``,
                          picks the fast native lane per page and only routes
                          scanned / math-heavy pages through PaddleOCR-VL.
      * ``"force_ocr"`` — ``native_fast_lane=False``: every page through
                          PaddleOCR-VL. Slower but robust to unreliable
                          text layers.

    ``free_translation_vram_first``: when True, evict any resident HuggingFace
    translation model from the GPU *before* spawning the PaddleOCR-VL worker.
    The OCR worker and a large (3B) translation backend do not fit together on
    smaller cards; freeing first makes them run sequentially and avoids the
    watchdog stall. Callers that translate afterwards should set this — the
    model is reloaded transparently on the next translate call.
    """
    # Deferred imports: pulling auto_ocr eagerly would drag the vision-model
    # subprocess wiring into every translate_engine import.
    from core import auto_ocr, doc_ir

    if free_translation_vram_first:
        # Give the OCR worker a clean GPU. Deferred import avoids a cycle at
        # module load (engine imports are otherwise lazy inside this package).
        from .engine import free_translation_vram

        free_translation_vram()

    ocr_prog = _ocr_progress_bridge(progress_cb, "OCR")

    with tempfile.TemporaryDirectory(prefix="tl_translate_ocr_") as tmp:
        input_path = Path(tmp) / source_name
        input_path.write_bytes(pdf_bytes)
        workspace = Path(tmp) / "ws"
        workspace.mkdir()

        document = auto_ocr.process_document(
            input_path,
            workspace,
            native_fast_lane=(pdf_type != "force_ocr"),
            progress=ocr_prog,
            source_name=source_name,
        )
        md_text = doc_ir.to_markdown(
            document, asset_dir="assets", embed_assets=True,
        )
        assets = doc_ir.collect_assets(document)
    return md_text, dict(assets)


def translate_pdf_to_markdown(
    pdf_bytes: bytes,
    translate_fn: TranslateFn,
    *,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
    pdf_type: str = "auto",
    source_name: str = "input.pdf",
) -> Tuple[str, Dict[str, bytes]]:
    """OCR the PDF into markdown, translate that markdown.

    Returns ``(translated_markdown, assets)``. Callers that want a single
    downloadable artifact can wrap the result with
    :func:`pack_markdown_bundle`.

    The OCR extraction frees any resident translation model from the GPU
    first, so the PaddleOCR-VL worker and the translation model run
    sequentially rather than competing for VRAM. The translation model is
    reloaded automatically for the :func:`translate_markdown` step below.
    """
    _report(progress_cb, 0, 1, "reading pdf")
    md_source, assets = pdf_to_markdown_bundle(
        pdf_bytes,
        pdf_type=pdf_type,
        source_name=source_name,
        progress_cb=progress_cb,
        free_translation_vram_first=True,
    )
    md_translated = translate_markdown(
        md_source,
        translate_fn,
        progress_cb=progress_cb,
        glossary=glossary,
    )
    return md_translated, assets


def pack_markdown_bundle(
    md_text: str,
    assets: Optional[Mapping[str, bytes]],
    *,
    stem: str,
) -> Tuple[bytes, str]:
    """Return ``(bytes, filename)`` — a ZIP when there are assets, else a ``.md``.

    ZIP layout:  ``<stem>.md``  +  ``assets/<name>.png`` per crop.
    """
    if not assets:
        return md_text.encode("utf-8"), f"{stem}.md"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{stem}.md", md_text)
        for name, data in assets.items():
            zf.writestr(f"assets/{name}", data)
    return buf.getvalue(), f"{stem}.md.zip"


def _split_by_ratio(text: str, ratios: List[float]) -> List[str]:
    """
    Split ``text`` into ``len(ratios)`` pieces whose lengths approximate the
    given ratios, snapping to nearest word boundaries.
    """
    n = len(ratios)
    if n <= 1:
        return [text]
    total = sum(ratios) or 1.0
    targets = [max(1, int(round(len(text) * r / total))) for r in ratios]
    diff = len(text) - sum(targets)
    targets[-1] += diff

    pieces: List[str] = []
    cursor = 0
    for tlen in targets[:-1]:
        end = cursor + tlen
        snap = text.find(" ", end)
        if snap == -1 or snap > cursor + tlen + 40:
            snap = end
        pieces.append(text[cursor:snap].strip())
        cursor = snap + 1 if snap < len(text) else len(text)
    pieces.append(text[cursor:].strip())
    return pieces


def _insert_autoshrink(
    page,
    bbox,
    text: str,
    preferred_size: float,
    font_path: Optional[str] = None,
) -> None:
    """
    Insert ``text`` into ``bbox`` on ``page`` with
    :meth:`fitz.Page.insert_textbox`, shrinking the font iteratively
    until it fits.

    If ``font_path`` points to a Unicode TTF that pymupdf can read, we
    use it (covers Greek/math/CJK/Arabic). Otherwise we fall back to
    the Base-14 Helvetica bundled with pymupdf.
    """
    fontsize = max(6.0, min(preferred_size, 14.0))

    if font_path:
        kw = dict(fontname="tl_uni", fontfile=font_path)
    else:
        kw = dict(fontname="helv")

    for _ in range(8):
        rc = page.insert_textbox(
            bbox,
            text,
            fontsize=fontsize,
            color=(0, 0, 0),
            align=0,
            expandtabs=4,
            **kw,
        )
        if rc >= 0:
            return
        fontsize *= 0.85
        if fontsize < 4.5:
            break
    # Final attempt: truncate with ellipsis rather than skip entirely.
    page.insert_textbox(
        bbox,
        text[: max(1, int(len(text) * 0.9))].rstrip() + "…",
        fontsize=max(4.5, fontsize),
        color=(0, 0, 0),
        align=0,
        expandtabs=4,
        **kw,
    )


# ===========================================================================
# XLSX (openpyxl)
# ===========================================================================


def translate_xlsx(
    xlsx_bytes: bytes,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
) -> bytes:
    """
    Translate a .xlsx file cell-by-cell, returning a new .xlsx file.

    Preserves:
    * Workbook / sheet structure, column widths, merged cells, styles.
    * Numeric, date, boolean, and formula cells (they are left untouched;
      formulas are never sent to the model).
    * Comments' contents are translated too.

    Limitations:
    * Charts and images embedded via drawings are preserved as-is (their
      textual labels are not translated).
    * Sheet *names* are not translated — Excel range references rely on
      them and translation risks breaking cross-sheet formulas.
    """
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError(
            "openpyxl is required to translate .xlsx files."
        ) from exc

    wb = load_workbook(io.BytesIO(xlsx_bytes))

    # First pass — count translatable strings for a smooth progress bar.
    translatable: list = []  # (cell, kind='cell' | 'comment')
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                val = cell.value
                if not isinstance(val, str):
                    continue
                # Skip formula strings (openpyxl gives data_type == 'f' or
                # the string starts with '=') and empty/whitespace cells.
                if getattr(cell, "data_type", None) == "f":
                    continue
                if val.startswith("="):
                    continue
                if not val.strip():
                    continue
                translatable.append((cell, "cell"))

                comment = getattr(cell, "comment", None)
                if comment is not None and comment.text and comment.text.strip():
                    translatable.append((cell, "comment"))

    total = max(1, len(translatable))
    texts = [
        cell.value if kind == "cell" else cell.comment.text
        for cell, kind in translatable
    ]
    _report(progress_cb, 0, total, "translating xlsx")
    translated_list = shielded_translate_many(texts, translate_fn, glossary=glossary)
    _report(progress_cb, total, total, "translating xlsx")

    for (cell, kind), translated in zip(translatable, translated_list):
        if kind == "cell":
            cell.value = translated
        else:
            comment = cell.comment
            comment.text = translated
            cell.comment = comment

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# ===========================================================================
# PPTX (python-pptx)
# ===========================================================================


def _pptx_paragraph_text(paragraph) -> str:
    """Join all run texts inside a python-pptx paragraph."""
    return "".join(r.text or "" for r in paragraph.runs)


def _pptx_rewrite_paragraph(paragraph, translated: str) -> None:
    """
    Replace the paragraph's visible text with ``translated``.

    Puts the whole translated string into the first run and empties the
    remaining runs, preserving the first run's font/color and the
    paragraph-level style. Intra-paragraph run styling is lost (same
    trade-off as the DOCX pipeline).
    """
    runs = list(paragraph.runs)
    if not runs:
        return
    runs[0].text = translated
    for r in runs[1:]:
        r.text = ""


def _iter_pptx_text_frames(prs):
    """Yield every ``text_frame`` from slides, notes, and table cells."""
    for slide in prs.slides:
        for shape in slide.shapes:
            yield from _iter_shape_text_frames(shape)
        if slide.has_notes_slide:
            notes = slide.notes_slide
            for shape in notes.shapes:
                yield from _iter_shape_text_frames(shape)


def _iter_shape_text_frames(shape):
    """Recurse into groups/tables/text-frames and yield each text frame."""
    # Grouped shapes.
    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        try:
            for sub in shape.shapes:
                yield from _iter_shape_text_frames(sub)
            return
        except AttributeError:
            pass
    # Tables: iterate cells.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame is not None:
                    yield cell.text_frame
        return
    # Regular text-bearing shapes.
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame


def translate_pptx(
    pptx_bytes: bytes,
    translate_fn: TranslateFn,
    progress_cb: ProgressCb = None,
    glossary: Glossary = None,
) -> bytes:
    """
    Translate a .pptx file, returning a new .pptx file.

    Preserves slide layout, shape positions, images, colors, tables,
    speaker notes, and paragraph-level style. Same intra-paragraph run
    styling caveat as :func:`translate_docx`.

    Raises ``RuntimeError`` if ``python-pptx`` is not available in the
    running Python environment.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required to translate .pptx files."
        ) from exc

    prs = Presentation(io.BytesIO(pptx_bytes))

    # Collect (text_frame, paragraph) pairs so we can report a total.
    paragraphs: list = []
    for tf in _iter_pptx_text_frames(prs):
        for para in tf.paragraphs:
            paragraphs.append(para)

    total = max(1, len(paragraphs))
    texts = [_pptx_paragraph_text(para) for para in paragraphs]
    _report(progress_cb, 0, total, "translating pptx")
    translated_list = shielded_translate_many(texts, translate_fn, glossary=glossary)
    _report(progress_cb, total, total, "translating pptx")

    for para, text, translated in zip(paragraphs, texts, translated_list):
        if not text.strip():
            continue
        _pptx_rewrite_paragraph(para, translated)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

